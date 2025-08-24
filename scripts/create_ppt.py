# -*- coding: utf-8 -*-
"""
Create a PowerPoint deck embedding the generated charts.
"""
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # Title
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_title_text(prs, title, body):
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body.split("\n")):
        if i == 0:
            tf.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0

def add_image_slide(prs, title, image_path):
    slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.5)
    height = Inches(5.0)
    slide.shapes.add_picture(str(image_path), left, top, height=height)

def build_presentation(output_dir: Path, out_pptx: Path, author: str = "Seu Nome"):
    prs = Presentation()

    # Slide 1 — Capa
    add_title_slide(
        prs,
        "Diagnóstico de Turnover & Ações Prioritárias",
        f"Período: Jan/2021 – Dez/2022\n{author} — {date.today().strftime('%d/%m/%Y')}"
    )

    # Slide 2 — Sumário
    add_title_text(
        prs, "Sumário",
        "Pilar 1 — Diagnóstico Estatístico\n"
        "Pilar 2 — Melhoria Contínua (DMAIC)\n"
        "Pilar 3 — Preditivo (ML)\n"
        "Pilar 4 — Plano de Ação"
    )

    # Slide 3 — Base e Método
    add_title_text(
        prs, "Base & Método",
        "• Base: 24 meses (Jan/21–Dez/22), hires/terminações/voluntário/involuntário, headcount total\n"
        "• Headcount médio: média entre total_{t-1} e total_t (primeiro mês usa total do mês)\n"
        "• KPI 1: Turnover (rotatividade) = [(Admissões + Desligamentos) / 2] / HC total\n"
        "• KPI 2: Taxa de desligamento = Desligamentos / HC médio\n"
        "• KPIs voluntário e involuntário = Voluntários / HC médio; Involuntários / HC médio"
    )

    # Slide 4 — Série Temporal (Desligamento)
    add_image_slide(prs, "Série: Taxa de Desligamento Mensal", output_dir / "taxa_desligamento.png")

    # Slide 5 — Série Temporal (Turnover Rotatividade)
    add_image_slide(prs, "Série: Turnover (Rotatividade) Mensal", output_dir / "turnover_rotatividade.png")

    # Slide 6 — Voluntário vs Involuntário
    add_image_slide(prs, "Turnover Voluntário", output_dir / "turnover_voluntario.png")
    add_image_slide(prs, "Turnover Involuntário", output_dir / "turnover_involuntario.png")

    # Slide 7 — Sazonalidade
    add_image_slide(prs, "Sazonalidade: Taxa de Desligamento", output_dir / "seasonality_desligamento.png")

    # Slide 8 — Tendência (Médias Móveis)
    add_image_slide(prs, "Tendência: MM3 Taxa de Desligamento", output_dir / "desligamento_ma3.png")
    add_image_slide(prs, "Tendência: MM3 Turnover (Rotatividade)", output_dir / "turnover_ma3.png")

    # Slide 9 — Próximos Passos (exemplo)
    add_title_text(
        prs, "Hipóteses e Próximos Passos",
        "H1 (Voluntário): compa-ratio e eNPS baixos aumentam saídas.\n"
        "H2 (Involuntário): absenteísmo/performance baixos elevam desligamentos.\n"
        "Coletar: remuneração, carreira, eNPS, tenure, engajamento, horas extras, performance, motivo da saída."
    )

    prs.save(str(out_pptx))

if __name__ == "__main__":
    outdir = Path(__file__).resolve().parents[1] / "output"
    out_pptx = outdir / "case_turnover_apresentacao.pptx"
    build_presentation(outdir, out_pptx, author="Mario L. O. Porto")
